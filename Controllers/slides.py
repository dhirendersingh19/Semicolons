
import collections
import collections.abc
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Inches
from Models import slides_res
from Utils import utils
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION
import random


class Slides:
    def __init__(self, path):
        self.path = path

    def create_ppt(self, date, script):
        filename = self.deck_title("Own Cloud - Review", date, script)
        self.summary(f"PPT/{filename}", "Summary")
        self.title_content(f"PPT/{filename}", "Task Status")
        self.add_test_case_deck(f"PPT/{filename}", "Test Cases")
        self.work_progress_graph(f"PPT/{filename}", "Team Progress")
        self.budget_deck(f"PPT/{filename}", "Project Budget")
        res = slides_res.SlidesRes(
            f"/{utils.BASE_ENDPOINT}/v1/download-ppt/{filename}")
        return res, 201

    def deck_title(self, ppt_title, date, script):

        # create presentation
        prs = Presentation(self.path)
        # add title slide
        deck_slide_layout = prs.slide_layouts[0]
        deck_slide = prs.slides.add_slide(deck_slide_layout)
        title = deck_slide.shapes.title
        subtitle = deck_slide.placeholders[1]
        title.text = ppt_title
        subtitle.text = date
        filename = f"SPRINT1-{date}.pptx"
        prs.save(f"PPT/{filename}")

        return filename

    def summary(self, filepath, ppt_title):
        prs = Presentation(filepath)

        content_slide_layout = prs.slide_layouts[3]
        content_slide = prs.slides.add_slide(content_slide_layout)
        title = content_slide.shapes.title
        title.text = ppt_title

        # Add a bullet point paragraph
        bullet_paragraph = content_slide.shapes.placeholders[1]
        bullet_paragraph.text = "Hi this is alston, as a project manager I will be managing thise team. Siddesh started working on developing APIS for file-upload feature, JIRA-ticket:1401. Siyashree is working oc creating the cypress test cases for owncloud, JIRA-ticket:1402. She is has completed the test cases for the login page and Dashboard and other feature test cases are pending. Diraj is creating UI components for owncloud. He has completed the login component,, JIRA-ticket:1403 and other components are pending. Dhirender is working on mongoDB api to create, delete, update and fetch users from the mongoDb database, JIRA-ticket:1404."

        prs.save(filepath)

    def title_content(self, filepath, ppt_title):
        prs = Presentation(filepath)

        content_slide_layout = prs.slide_layouts[3]
        content_slide = prs.slides.add_slide(content_slide_layout)
        title = content_slide.shapes.title
        title.text = ppt_title

        # Add a bullet point paragraph
        bullet_paragraph = content_slide.shapes.placeholders[1].text_frame.add_paragraph(
        )
        bullet_paragraph.text = "Completed"
        bullet_paragraph.font.bold = True
        bullet_paragraph.level = 0

        completed_task = [
            "Alston has done most of the development work related to add cart functionality and identified/fixed a defect related to login functionality.",
            "Bob has completed the task of adding tests for the signup page and identified/fixed a defect related to OTP during signup.",
            "Jaz has created a database for order details and the backend team has started using it."
        ]
        for item in completed_task:
            self.display_points(content_slide.shapes, item)

        # Add a bullet point paragraph
        bullet_paragraph = content_slide.shapes.placeholders[1].text_frame.add_paragraph(
        )
        bullet_paragraph.text = "In-Progress"
        bullet_paragraph.font.bold = True
        bullet_paragraph.level = 0

        inprogress = [
            "Alston is raising a PR for the work done on Jira 1412 and looking into defects related to Jira 3452.",
            "Bob is working on adding test cases for the login page and is currently blocked due to issues in fetching data from the mock database.",
            "Jaz is currently resolving the defect related to frequent DB connection failures on Jira 3232."
        ]
        for item in inprogress:
            self.display_points(content_slide.shapes, item)

        prs.save(filepath)



    def add_test_case_deck(self, filepath, ppt_title):
        # create presentation object
        prs = Presentation(filepath)

        content_slide_layout = prs.slide_layouts[3]
        content_slide = prs.slides.add_slide(content_slide_layout)
        title = content_slide.shapes.title
        title.text = ppt_title

        subtitle = content_slide.placeholders[1]
        sp = subtitle.element
        sp.getparent().remove(sp)

        testcase = [{
            "TestcaseName": "Test Case Name",
            "Status": "Status",
            "Executed By": "Executed By",
        },
            {
            "TestcaseName": "Test Cases: Developing APIs for File-Upload Feature",
            "Status": "Pending",
            "Executed By": "Siddesh",
        },
        {
            "TestcaseName": "Creating Cypress Test Cases for OwnCloud, Status: Completed",
            "Status": "Completed",
            "Executed By": "Siyashree",
        },
        {
            "TestcaseName": "Test Cases: Creating UI Components for OwnCloud",
            "Status": "Pending",
            "Executed By": "Diraj",
        }
        ]

        # define table dimensions
        rows = len(testcase)
        cols = 3

        # define table coordinates and dimensions
        left = Inches(0.5)
        top = Inches(1)
        width = Inches(12)
        height = Inches(1)

        # add table to slide
        table = content_slide.shapes.add_table(
            rows, cols, left, top, width, height).table

        # populate table cells
        for i in range(0, rows):
            values = self.getTestCaseValue(testcase[i])
            for j in range(cols):
                cell = table.cell(i, j)
                cell.text = f"{values[j]}"

        # save presentation
        prs.save(filepath)

    def display_points(self, shapes, item):
        bullet_paragraph = shapes.placeholders[1].text_frame.add_paragraph()
        bullet_paragraph.text = item
        bullet_paragraph.font.bold = False
        bullet_paragraph.level = 1

    def getTestCaseValue(self, testcase):
        list = []
        list.append(testcase["TestcaseName"])
        list.append(testcase["Status"])
        list.append(testcase["Executed By"])

        return list

    def work_progress_graph(self, filepath, ppt_title):

        # create presentation object
        prs = Presentation(filepath)

        # add slide to presentation
        slide = prs.slides.add_slide(prs.slide_layouts[3])

        title = slide.shapes.title
        title.text = ppt_title

        subtitle = slide.placeholders[1]
        sp = subtitle.element
        sp.getparent().remove(sp)

        # define chart dimensions and coordinates
        left = Inches(0.5)
        top = Inches(1)
        width = Inches(12)
        height = Inches(6)

        # add chart to slide
        chart_data = CategoryChartData()
        chart_data.categories = ["SPRINT1", "SPRINT2", "SPRINT3"]
        chart_data.add_series("ALSTON", ("1", "4", "3"))
        chart_data.add_series("SIDDESH", ("3", "5", "2"))
        chart_data.add_series("DHIRAJ", ("2", "3", "4"))
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, left, top, width, height, chart_data).chart

        # # set series colors
        plot = chart.plots[0]

        for i in range(2):
            r = random.randrange(0, 255)
            g = random.randrange(0, 255)
            b = random.randrange(0, 255)
            plot.series[i].format.line.color.rgb = RGBColor(r, g, b)

        # save presentation
        prs.save(filepath)

    def budget_deck(self, filepath, ppt_title):
        # create a new PowerPoint presentation
        prs = Presentation(filepath)

        # add a new slide to the presentation
        slide = prs.slides.add_slide(prs.slide_layouts[3])

        title = slide.shapes.title
        title.text = ppt_title

        subtitle = slide.placeholders[1]
        sp = subtitle.element
        sp.getparent().remove(sp)

        # define the chart data
        chart_data = CategoryChartData()
        chart_data.categories = ['Sprint 1', 'Sprint 2', 'Sprint 3']
        chart_data.add_series('Cloud Budget', (2000, 1000, 2500))
        chart_data.add_series('Resource Budget', (1000, 500, 2000))

        # add a chart to the slide
        x, y, cx, cy = Inches(0.5), Inches(1), Inches(12), Inches(6)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart

        # set the chart title
        chart.chart_title.text_frame.text = 'Budget'

        # set the legend position
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        chart.value_axis.axis_title.text_frame.text = 'USD'

        for i in range(2):
            # set the data labels position and number format for series 1
            data_labels_1 = chart.series[i].data_labels
            data_labels_1.number_format = '0.0'  # set label format
            data_labels_1.position = XL_LABEL_POSITION.OUTSIDE_END  # set label position

            # loop through each point in series 1 to set the label text
            for idx, point in enumerate(chart.series[i].points):
                label = point.data_label
                # set label text
                label.text = chart_data._series[i].values[idx]

        # save the PowerPoint presentation
        prs.save(filepath)
